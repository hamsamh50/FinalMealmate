"""
Generates a presentation report for the MealMate project in both:
  - PowerPoint (MealMate_Report.pptx)
  - PDF        (MealMate_Report.pdf)

Run:  python generate_report.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle,
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT


BASE_DIR = Path(__file__).resolve().parent

PROJECT_TITLE = "MealMate"
PROJECT_SUBTITLE = "Online Food Ordering & Delivery System"
PROJECT_TAGLINE = "A Django-based web application with Razorpay payment integration"

# ---------------------------------------------------------------------------
# Slide / section content
# ---------------------------------------------------------------------------
SECTIONS = [
    {
        "title": "Introduction",
        "bullets": [
            "MealMate is an online food ordering and delivery web application.",
            "Built using the Django web framework (Python) with SQLite as the database.",
            "Supports two roles: Admin (manages restaurants & menus) and Customer (browses, orders, pays).",
            "Integrated with the Razorpay payment gateway for secure online payments.",
        ],
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Customers need a simple, centralized platform to browse multiple restaurants and order food.",
            "Restaurants need a manageable way to maintain their menus and prices online.",
            "Manual ordering processes (phone / in-person) are slow and error-prone.",
            "Need for a secure online payment workflow with order confirmation.",
        ],
    },
    {
        "title": "Objectives",
        "bullets": [
            "Provide user authentication (sign up / sign in) for customers and admins.",
            "Enable admins to add, update, and delete restaurants and menu items.",
            "Allow customers to browse restaurants, view menus, and add items to a cart.",
            "Facilitate secure checkout via Razorpay and display an order confirmation.",
        ],
    },
    {
        "title": "Technology Stack",
        "bullets": [
            "Backend: Python 3, Django 5.1.14",
            "Frontend: HTML5, CSS3, Django Template Language",
            "Database: SQLite (default Django DB)",
            "Payment Gateway: Razorpay (razorpay 2.0.1 SDK)",
            "Other Tools: Django Admin, Django ORM, Django URL routing",
        ],
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Client (Browser) <-> Django Views <-> Django ORM <-> SQLite Database",
            "Templates render dynamic HTML pages using context data.",
            "Razorpay client created in checkout view to generate an order ID.",
            "Frontend payment widget collects payment, then redirects to orders page.",
            "MVC-style separation: Models (data), Views (logic), Templates (presentation).",
        ],
    },
    {
        "title": "Modules",
        "bullets": [
            "Authentication: signup, signin, role-based redirection (admin vs. customer).",
            "Restaurant Management (Admin): add / update / delete restaurants.",
            "Menu Management (Admin): add menu items linked to a restaurant.",
            "Customer Browsing: view restaurant list and per-restaurant menu.",
            "Cart: add items to cart, view cart, compute total price.",
            "Checkout & Payment: Razorpay order creation & confirmation page.",
        ],
    },
    {
        "title": "Database Models",
        "bullets": [
            "Customer: username, password, email, mobile, address.",
            "Restaurant: name, picture (URL), cuisine, rating.",
            "Item: restaurant (FK), name, description, price, vegetarian flag, picture.",
            "Cart: customer (FK), items (M2M to Item), total_price() helper.",
            "Relationships: Restaurant 1<->N Item ; Customer 1<->1 Cart ; Cart N<->N Item.",
        ],
    },
    {
        "title": "Key Features",
        "bullets": [
            "Two-tier login flow: 'admin' user routes to admin dashboard, others to customer home.",
            "Restaurant catalog with images, cuisine type, and rating.",
            "Per-restaurant menu with vegetarian / non-vegetarian flag.",
            "Add-to-cart, cart total computation, and clear-cart-on-order behavior.",
            "Razorpay-powered online payments in INR (test mode keys configured).",
        ],
    },
    {
        "title": "URL Routing Highlights",
        "bullets": [
            "/                       Landing page",
            "/open_signin, /signin   Authentication",
            "/open_signup, /signup   Customer registration",
            "/open_add_restaurant    Admin: add restaurant page",
            "/view_menu/<id>/<user>  Customer: view a restaurant's menu",
            "/show_cart/<user>       Customer: cart page",
            "/checkout/<user>/       Razorpay checkout",
            "/orders/<user>/         Order confirmation",
        ],
    },
    {
        "title": "Workflow - Customer",
        "bullets": [
            "1. Sign up / Sign in.",
            "2. View list of restaurants on the home page.",
            "3. Open a restaurant's menu and add items to cart.",
            "4. Open cart, review items and total price.",
            "5. Proceed to checkout (Razorpay payment).",
            "6. On success, view orders confirmation; cart is cleared.",
        ],
    },
    {
        "title": "Workflow - Admin",
        "bullets": [
            "1. Sign in as 'admin'.",
            "2. Add a new restaurant (name, picture URL, cuisine, rating).",
            "3. View all restaurants and update or delete any of them.",
            "4. Update menu of a restaurant by adding new items.",
            "5. Items are linked to their parent restaurant via a foreign key.",
        ],
    },
    {
        "title": "Advantages",
        "bullets": [
            "Centralized platform for multiple restaurants and customers.",
            "Quick and convenient food ordering experience.",
            "Secure online payment via a trusted gateway (Razorpay).",
            "Simple admin interface to keep restaurant data up to date.",
            "Easy to extend (new fields, new modules) thanks to Django's ORM.",
        ],
    },
    {
        "title": "Limitations",
        "bullets": [
            "Passwords are stored as plain text (CharField) - not production safe.",
            "No real authentication / session management built on Django auth.",
            "No order history persistence (cart is cleared on confirmation).",
            "No delivery tracking, ratings, or reviews from customers yet.",
            "Static media only via image URLs - no file uploads.",
        ],
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Use Django's built-in auth system with hashed passwords.",
            "Persist orders in a dedicated Order / OrderItem model with status tracking.",
            "Add delivery agent module and live order tracking.",
            "Allow customers to rate and review restaurants and items.",
            "Mobile-responsive UI and a REST API for a future mobile app.",
            "Deploy to a cloud platform (PythonAnywhere / Render / AWS).",
        ],
    },
    {
        "title": "Conclusion",
        "bullets": [
            "MealMate demonstrates a complete food-ordering workflow on Django.",
            "Covers user roles, CRUD operations, cart management, and online payments.",
            "Provides a solid foundation for scaling into a production-grade system.",
            "Showcases practical use of Django ORM, templates, and third-party integrations.",
        ],
    },
]


# ---------------------------------------------------------------------------
# PowerPoint generation
# ---------------------------------------------------------------------------
COLOR_PRIMARY = RGBColor(0xC0, 0x39, 0x2B)   # red-ish (food theme)
COLOR_SECONDARY = RGBColor(0x2C, 0x3E, 0x50)  # dark slate
COLOR_ACCENT = RGBColor(0xF3, 0x9C, 0x12)     # warm orange
COLOR_LIGHT_BG = RGBColor(0xFD, 0xF6, 0xEC)   # cream
COLOR_TEXT = RGBColor(0x2B, 0x2B, 0x2B)


def _add_background(slide, prs, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send background to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_accent_bar(slide, prs):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()


def build_pptx(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # ---- Title slide ----
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, prs, COLOR_LIGHT_BG)

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.6), prs.slide_width, Inches(2.3)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    title_tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(12), Inches(1.2))
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = PROJECT_TITLE
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sub_tb = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(12), Inches(0.7))
    tf = sub_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = PROJECT_SUBTITLE
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tag_tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(0.6))
    tf = tag_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = PROJECT_TAGLINE
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.color.rgb = COLOR_SECONDARY

    foot_tb = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12), Inches(0.5))
    tf = foot_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = "Project Presentation Report"
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_SECONDARY

    # ---- Content slides ----
    for idx, section in enumerate(SECTIONS, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _add_background(slide, prs, RGBColor(0xFF, 0xFF, 0xFF))
        _add_accent_bar(slide, prs)

        # Title
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section["title"]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = COLOR_SECONDARY

        # Underline accent
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.45), Inches(1.6), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()

        # Bullets
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(5.0))
        tf = body.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(section["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(20)
            run.font.color.rgb = COLOR_TEXT

        # Footer with slide number
        foot = slide.shapes.add_textbox(Inches(11.2), Inches(7.05), Inches(2), Inches(0.4))
        tf = foot.text_frame
        p = tf.paragraphs[0]
        p.alignment = 2
        run = p.add_run()
        run.text = f"{PROJECT_TITLE}  |  Slide {idx + 1}"
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_SECONDARY

    # ---- Thank-you slide ----
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, prs, COLOR_PRIMARY)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(80)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = "Questions & Discussion"
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(str(output_path))
    print(f"Saved PPT  -> {output_path}")


# ---------------------------------------------------------------------------
# PDF generation
# ---------------------------------------------------------------------------
def build_pdf(output_path: Path):
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        leftMargin=0.7 * inch,
        rightMargin=0.7 * inch,
        topMargin=0.8 * inch,
        bottomMargin=0.8 * inch,
        title=f"{PROJECT_TITLE} - Project Report",
        author="MealMate",
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "ReportTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=34,
        leading=40,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#C0392B"),
        spaceAfter=10,
    )
    subtitle_style = ParagraphStyle(
        "ReportSubtitle",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=18,
        leading=22,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#2C3E50"),
        spaceAfter=6,
    )
    tagline_style = ParagraphStyle(
        "Tagline",
        parent=styles["Italic"],
        fontSize=12,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#555555"),
        spaceAfter=24,
    )
    h1 = ParagraphStyle(
        "H1",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#2C3E50"),
        spaceBefore=4,
        spaceAfter=10,
    )
    bullet_style = ParagraphStyle(
        "Bullet",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=11.5,
        leading=16,
        leftIndent=18,
        bulletIndent=4,
        alignment=TA_LEFT,
        textColor=colors.HexColor("#2B2B2B"),
        spaceAfter=4,
    )
    label_style = ParagraphStyle(
        "Label",
        parent=styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=12,
        textColor=colors.HexColor("#C0392B"),
    )

    story = []

    # ---- Cover ----
    story.append(Spacer(1, 1.4 * inch))
    story.append(Paragraph(PROJECT_TITLE, title_style))
    story.append(Paragraph(PROJECT_SUBTITLE, subtitle_style))
    story.append(Paragraph(PROJECT_TAGLINE, tagline_style))

    overview_data = [
        ["Project Name", PROJECT_TITLE],
        ["Type", "Web Application (Django)"],
        ["Domain", "Online Food Ordering & Delivery"],
        ["Database", "SQLite"],
        ["Payment Gateway", "Razorpay"],
        ["Document", "Project Presentation Report"],
    ]
    tbl = Table(overview_data, colWidths=[1.7 * inch, 4.0 * inch])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#C0392B")),
        ("TEXTCOLOR", (0, 0), (0, -1), colors.white),
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("BACKGROUND", (1, 0), (1, -1), colors.HexColor("#FDF6EC")),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#DDDDDD")),
    ]))
    story.append(Spacer(1, 0.3 * inch))
    story.append(tbl)

    story.append(PageBreak())

    # ---- Table of contents ----
    story.append(Paragraph("Table of Contents", h1))
    toc_items = [[str(i + 1), s["title"]] for i, s in enumerate(SECTIONS)]
    toc_tbl = Table(toc_items, colWidths=[0.5 * inch, 5.5 * inch])
    toc_tbl.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 12),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("TEXTCOLOR", (0, 0), (0, -1), colors.HexColor("#C0392B")),
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("LINEBELOW", (0, 0), (-1, -1), 0.3, colors.HexColor("#EEEEEE")),
    ]))
    story.append(toc_tbl)
    story.append(PageBreak())

    # ---- Sections ----
    for i, section in enumerate(SECTIONS, start=1):
        story.append(Paragraph(f"{i}. {section['title']}", h1))
        for bullet in section["bullets"]:
            story.append(Paragraph(bullet, bullet_style, bulletText="•"))
        story.append(Spacer(1, 0.25 * inch))

    # ---- Closing ----
    story.append(PageBreak())
    story.append(Spacer(1, 2.5 * inch))
    story.append(Paragraph("Thank You", title_style))
    story.append(Paragraph("Questions & Discussion", subtitle_style))

    def _on_page(canvas, doc_):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#C0392B"))
        canvas.setLineWidth(2)
        canvas.line(
            doc_.leftMargin, A4[1] - 0.55 * inch,
            A4[0] - doc_.rightMargin, A4[1] - 0.55 * inch,
        )
        canvas.setFont("Helvetica-Bold", 9)
        canvas.setFillColor(colors.HexColor("#C0392B"))
        canvas.drawString(doc_.leftMargin, A4[1] - 0.45 * inch, PROJECT_TITLE)
        canvas.setFont("Helvetica", 9)
        canvas.setFillColor(colors.HexColor("#555555"))
        canvas.drawRightString(
            A4[0] - doc_.rightMargin, A4[1] - 0.45 * inch,
            "Project Presentation Report",
        )
        canvas.setFont("Helvetica", 9)
        canvas.drawCentredString(
            A4[0] / 2.0, 0.4 * inch, f"Page {doc_.page}"
        )
        canvas.restoreState()

    doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    print(f"Saved PDF  -> {output_path}")


# ---------------------------------------------------------------------------
if __name__ == "__main__":
    pptx_path = BASE_DIR / "MealMate_Report.pptx"
    pdf_path = BASE_DIR / "MealMate_Report.pdf"
    build_pptx(pptx_path)
    build_pdf(pdf_path)
    print("Done.")
